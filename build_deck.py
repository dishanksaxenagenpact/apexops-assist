"""ApexAssist · Leadership Deck (3 slides, 16:9, Genpact-native)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

INK       = RGBColor(0x0A, 0x0F, 0x1F)   # deep navy
INK_2     = RGBColor(0x12, 0x18, 0x2C)
INK_3     = RGBColor(0x1B, 0x22, 0x3A)
TX_1      = RGBColor(0xFF, 0xFF, 0xFF)
TX_2      = RGBColor(0xB8, 0xC1, 0xDE)
TX_3      = RGBColor(0x7F, 0x8A, 0xAE)
CORAL     = RGBColor(0xFF, 0x4F, 0x59)
VIOLET    = RGBColor(0x97, 0x75, 0xFA)
MAGENTA   = RGBColor(0xF0, 0x62, 0x92)
CYAN      = RGBColor(0x22, 0xD3, 0xEE)
GREEN     = RGBColor(0x4A, 0xDE, 0x80)
AMBER     = RGBColor(0xFB, 0xBF, 0x24)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

blank = prs.slide_layouts[6]

def bg(slide, color=INK):
    """Solid background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def rect(slide, x, y, w, h, fill=INK_2, line=None, line_w=0.75, rounded=False, corner=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = corner
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def txt(slide, x, y, w, h, text, size=14, color=TX_1, bold=False, font='Calibri',
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def pill(slide, x, y, w, h, label, fill=INK_3, color=TX_2, size=9, font='Consolas', border=None):
    s = rect(slide, x, y, w, h, fill=fill, line=border, line_w=0.5, rounded=True, corner=0.5)
    tb = txt(slide, x, y, w, h, label, size=size, color=color, font=font,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s

def gradient_band(slide, x, y, w, h, c1=VIOLET, c2=MAGENTA, vertical=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    # XML gradient
    sppr = shp.fill._xPr.spPr
    # remove default fill
    for tag in ('a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:noFill'):
        for el in sppr.findall(qn(tag)):
            sppr.remove(el)
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    grad_xml = f'''<a:gradFill rotWithShape="1" xmlns:a="{ns}">
        <a:gsLst>
          <a:gs pos="0"><a:srgbClr val="{c1:02X}{c1:02X}{c1:02X}"/></a:gs>
          <a:gs pos="100000"><a:srgbClr val="{c2:02X}{c2:02X}{c2:02X}"/></a:gs>
        </a:gsLst>
        <a:lin ang="{'5400000' if vertical else '0'}" scaled="0"/>
      </a:gradFill>'''
    # Simpler: just use solid color for safety
    shp.fill.solid(); shp.fill.fore_color.rgb = c1
    return shp

def accent_bar(slide, x, y, w=0.08, h=0.45, color=CORAL):
    return rect(slide, x, y, w, h, fill=color)

# =====================================================
# SLIDE 1 — Title / Vision
# =====================================================
s1 = prs.slides.add_slide(blank)
bg(s1, INK)

# Decorative aurora bands (rectangles with low-opacity not natively supported — use color overlays)
rect(s1, 8.0, -1.0, 6.5, 5.5, fill=RGBColor(0x1A, 0x12, 0x35))   # violet glow
rect(s1, 10.5, 1.5, 4.5, 4.5, fill=RGBColor(0x2A, 0x10, 0x28))   # magenta glow
rect(s1, -1, 4.5, 8, 4, fill=RGBColor(0x0F, 0x18, 0x2F))         # bottom dark

# Coral accent bar (brand)
rect(s1, 0, 0, 13.333, 0.12, fill=CORAL)

# Top meta row
pill(s1, 0.55, 0.45, 2.6, 0.34, "● GENPACT · GENAI HACKATHON 2026", fill=INK_3, color=CORAL, size=9)
pill(s1, 3.30, 0.45, 1.9, 0.34, "POWERED BY CLAUDE", fill=INK_3, color=VIOLET, size=9)

# Brand mark — hex
hex_box = s1.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.55), Inches(1.15), Inches(0.85), Inches(0.95))
hex_box.fill.solid(); hex_box.fill.fore_color.rgb = INK_3
hex_box.line.color.rgb = VIOLET
hex_box.line.width = Pt(1.5)
txt(s1, 0.55, 1.15, 0.85, 0.95, "⬡", size=28, color=VIOLET, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

txt(s1, 1.6, 1.20, 5, 0.45, "ApexAssist", size=14, color=TX_2, bold=True, font='Calibri')
txt(s1, 1.6, 1.50, 5, 0.4, "by Team ApexOps", size=11, color=TX_3, font='Consolas')

# Hero title
txt(s1, 0.55, 2.30, 12, 1.4, "Autonomous ITOps,", size=64, color=TX_1, bold=True, font='Calibri')
txt(s1, 0.55, 3.20, 12, 1.4, "Orchestrated by AI.", size=64, color=VIOLET, bold=True, font='Calibri')

# Sub headline
txt(s1, 0.55, 4.55, 11, 0.6,
    "One Claude-powered Orchestrator. Seven specialist agents. Zero swivel-chair.",
    size=20, color=TX_2, font='Calibri')

# Three pillars row
pillar_y = 5.55
pillar_w = 4.0
pillar_h = 1.35
pillar_data = [
    ("⚡  −72%", "Mean Time to Resolve", VIOLET),
    ("🛡  6×", "Built-in Safety Guardrails", CORAL),
    ("💰  $3.2M", "Annual Savings / 100 inc / mo", CYAN),
]
for i, (head, sub, c) in enumerate(pillar_data):
    px = 0.55 + i*(pillar_w + 0.30)
    rect(s1, px, pillar_y, pillar_w, pillar_h, fill=INK_2, line=INK_3, rounded=True, corner=0.08)
    accent_bar(s1, px, pillar_y + 0.18, w=0.08, h=0.95, color=c)
    txt(s1, px+0.30, pillar_y+0.20, pillar_w-0.4, 0.55, head, size=24, color=c, bold=True)
    txt(s1, px+0.30, pillar_y+0.78, pillar_w-0.4, 0.5, sub, size=12, color=TX_2)

# Footer
txt(s1, 0.55, 7.10, 12, 0.3,
    "Dishank Saxena · Sreenivasan Masanamuthu · Mihir Patel",
    size=10, color=TX_3, font='Consolas')

# =====================================================
# SLIDE 2 — How It Works (architecture + agents)
# =====================================================
s2 = prs.slides.add_slide(blank)
bg(s2, INK)
rect(s2, 0, 0, 13.333, 0.12, fill=VIOLET)

pill(s2, 0.55, 0.30, 2.0, 0.32, "▸ HOW IT WORKS", fill=INK_3, color=VIOLET, size=9)
txt(s2, 0.55, 0.75, 12, 0.6, "Plan → Act → Reflect.", size=30, color=TX_1, bold=True)
txt(s2, 0.55, 1.30, 12, 0.5,
    "An Orchestrator delegates to 7 specialist agents — each with a JSON contract, audited tool-use, and SLO gates.",
    size=13, color=TX_2)

# Center hub
HUB_X, HUB_Y, HUB_R = 6.40, 4.20, 0.85
hex_hub = s2.shapes.add_shape(MSO_SHAPE.HEXAGON,
    Inches(HUB_X), Inches(HUB_Y), Inches(HUB_R*2), Inches(HUB_R*1.7))
hex_hub.fill.solid(); hex_hub.fill.fore_color.rgb = INK_3
hex_hub.line.color.rgb = VIOLET; hex_hub.line.width = Pt(2.5)
txt(s2, HUB_X, HUB_Y+0.15, HUB_R*2, 0.6, "⬡", size=42, color=VIOLET, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s2, HUB_X, HUB_Y+0.70, HUB_R*2, 0.4, "ORCHESTRATOR", size=11, color=TX_1, bold=True,
    align=PP_ALIGN.CENTER, font='Consolas')
txt(s2, HUB_X, HUB_Y+1.05, HUB_R*2, 0.35, "Claude Sonnet 4", size=9, color=TX_2,
    align=PP_ALIGN.CENTER, font='Consolas')

# 7 agents arranged in a circle around hub
import math
agents = [
    ("🔍", "Monitoring",   "signal correlation",  AMBER),
    ("📋", "Triage",       "severity scoring",    MAGENTA),
    ("🧬", "RCA",          "causal inference",    VIOLET),
    ("🔧", "Self Heal",    "auto-remediation",    GREEN),
    ("✅", "Validator",    "SLO gate",            MAGENTA),
    ("📢", "Smart Notify", "multi-channel",       CYAN),
    ("📚", "Learning",     "continuous eval",     VIOLET),
]
CX, CY = HUB_X + HUB_R, HUB_Y + HUB_R*0.85
R_OUT = 2.45
N = len(agents)
for i,(emoji,name,sub,clr) in enumerate(agents):
    ang = -math.pi/2 + (i/N)*2*math.pi
    ax = CX + R_OUT*math.cos(ang) - 0.85
    ay = CY + R_OUT*math.sin(ang)*0.85 - 0.45
    w, h = 1.7, 0.9
    rect(s2, ax, ay, w, h, fill=INK_2, line=clr, line_w=1.0, rounded=True, corner=0.18)
    txt(s2, ax+0.05, ay+0.06, 0.4, 0.4, emoji, size=14,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s2, ax+0.40, ay+0.05, w-0.45, 0.38, name, size=11, color=TX_1, bold=True)
    txt(s2, ax+0.40, ay+0.38, w-0.45, 0.40, sub, size=8, color=clr, font='Consolas')

# Left rail — Plan-Act-Reflect loop callouts
rail_x, rail_y = 0.55, 2.30
loop_items = [
    ("PLAN",    "Decompose the incident, pick agents & tools", VIOLET),
    ("ACT",     "Invoke tools via OPA-gated catalogue",        CORAL),
    ("REFLECT", "Self-critique · escalate to human if needed", CYAN),
]
for i,(t,d,c) in enumerate(loop_items):
    y = rail_y + i*1.05
    rect(s2, rail_x, y, 2.95, 0.92, fill=INK_2, rounded=True, corner=0.12)
    accent_bar(s2, rail_x, y+0.12, 0.07, 0.68, c)
    txt(s2, rail_x+0.25, y+0.10, 2.6, 0.38, t, size=11, color=c, bold=True, font='Consolas')
    txt(s2, rail_x+0.25, y+0.40, 2.6, 0.5, d, size=10, color=TX_2)

# Right rail — Guardrails
rg_x = 10.10
rg_items = [
    ("Human gates",       "for risky tool calls"),
    ("OPA policy",        "blast-radius allowlists"),
    ("PII redaction",     "ingress + egress"),
    ("Immutable audit",   "every prompt & action"),
]
txt(s2, rg_x, 2.0, 2.95, 0.4, "RESPONSIBLE AI", size=10, color=CORAL, bold=True, font='Consolas')
for i,(t,d) in enumerate(rg_items):
    y = 2.45 + i*0.85
    rect(s2, rg_x, y, 2.75, 0.74, fill=INK_2, rounded=True, corner=0.15)
    accent_bar(s2, rg_x, y+0.10, 0.07, 0.54, CORAL)
    txt(s2, rg_x+0.22, y+0.06, 2.5, 0.34, t, size=11, color=TX_1, bold=True)
    txt(s2, rg_x+0.22, y+0.36, 2.5, 0.36, d, size=9, color=TX_2, font='Consolas')

# Bottom strip
txt(s2, 0.55, 7.15, 12, 0.3,
    "Memory: pgvector + Neo4j  ·  Event Bus: Kafka  ·  Action: Argo Workflows  ·  Evals: PromptLayer",
    size=10, color=TX_3, font='Consolas', align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 3 — Impact & Why we win
# =====================================================
s3 = prs.slides.add_slide(blank)
bg(s3, INK)
rect(s3, 0, 0, 13.333, 0.12, fill=CORAL)

pill(s3, 0.55, 0.30, 2.0, 0.32, "▸ THE PAYOFF", fill=INK_3, color=CORAL, size=9)
txt(s3, 0.55, 0.75, 12, 0.7, "Strategic, measurable, ready.", size=30, color=TX_1, bold=True)
txt(s3, 0.55, 1.30, 12, 0.5,
    "ApexAssist directly attacks the largest cost center in Genpact's ITS portfolio — incident toil.",
    size=13, color=TX_2)

# 4 KPIs row
kpis = [
    ("87%",   "MTTR reduction",         VIOLET),
    ("$3.2M", "Annual savings",         CORAL),
    ("80%",   "On-call toil reduction", CYAN),
    ("4 min", "Alert → all-clear",      GREEN),
]
kpi_y = 2.10
kw = 2.95
for i,(v,l,c) in enumerate(kpis):
    kx = 0.55 + i*(kw + 0.18)
    rect(s3, kx, kpi_y, kw, 1.55, fill=INK_2, rounded=True, corner=0.08)
    accent_bar(s3, kx, kpi_y+0.20, 0.08, 1.15, c)
    txt(s3, kx+0.28, kpi_y+0.18, kw-0.4, 0.85, v, size=36, color=c, bold=True)
    txt(s3, kx+0.28, kpi_y+1.00, kw-0.4, 0.45, l, size=11, color=TX_2)

# Rubric / why this wins
txt(s3, 0.55, 3.95, 12, 0.4, "WHY THIS WINS", size=11, color=CORAL, bold=True, font='Consolas')

cards = [
    ("🎯  Strategic for Genpact",
     "Largest IT services cost center. Proven economics. Demoable in 4 minutes."),
    ("🧠  True agentic system",
     "Not a chat wrapper — plan/act/reflect, JSON contracts, tool-use, memory, learning."),
    ("🛡️  Safe by construction",
     "OPA policy, human gates, PII redaction, audit log built-in from day one."),
]
cy = 4.35
cw = 4.0
ch = 1.85
for i,(t,d) in enumerate(cards):
    cx = 0.55 + i*(cw + 0.20)
    rect(s3, cx, cy, cw, ch, fill=INK_2, line=INK_3, rounded=True, corner=0.10)
    accent_bar(s3, cx, cy+0.25, 0.08, 1.35, [VIOLET, CORAL, CYAN][i])
    txt(s3, cx+0.30, cy+0.22, cw-0.45, 0.5, t, size=14, color=TX_1, bold=True)
    txt(s3, cx+0.30, cy+0.78, cw-0.45, ch-0.85, d, size=11, color=TX_2, spacing=1.25)

# CTA strip
rect(s3, 0.55, 6.60, 12.20, 0.65, fill=INK_2, line=VIOLET, line_w=1.2, rounded=True, corner=0.20)
txt(s3, 0.85, 6.60, 6.5, 0.65, "▶  Live demo",
    size=14, color=TX_1, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s3, 0.85, 6.60, 11.6, 0.65,
    "dishanksaxenagenpact.github.io/apexops-assist",
    size=12, color=VIOLET, font='Consolas',
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Footer
txt(s3, 0.55, 7.25, 12, 0.3,
    "Team ApexOps  ·  Dishank Saxena · Sreenivasan Masanamuthu · Mihir Patel  ·  Genpact GenAI Hackathon 2026",
    size=9, color=TX_3, font='Consolas')

out = '/Users/H500382/Downloads/Latest_Bitbucket/ApexAssist/ApexAssist_Leadership_Deck.pptx'
prs.save(out)
print('OK', out, os.path.getsize(out))
